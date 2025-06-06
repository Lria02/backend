# backend/extract/extractor.py

import fitz  # PyMuPDF
from pptx import Presentation
import tempfile
import os

def extract_text_from_file(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    text_data = []

    if ext == '.pdf':
        doc = fitz.open(file_path)
        for page in doc:
            text_data.append(page.get_text())

    elif ext == '.pptx':
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_data.append(shape.text)

    return "\n".join(text_data)
