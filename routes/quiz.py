from flask import Blueprint, request, jsonify
import requests
import os
from dotenv import load_dotenv
import tempfile
import re
from pptx import Presentation
import fitz  # PyMuPDF

load_dotenv()

quiz_bp = Blueprint('quiz', __name__)

def extract_text(file_path):
    all_text = []
    if file_path.lower().endswith(".pptx"):
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        all_text.append(text)
    elif file_path.lower().endswith(".pdf"):
        doc = fitz.open(file_path)
        for page in doc:
            text = page.get_text().strip()
            if text:
                all_text.append(text)
    return "\n".join(all_text)

@quiz_bp.route('/', methods=['POST'])
def generate_quiz():
    if 'file' not in request.files:
        return jsonify({"error": "No file uploaded"}), 400

    file = request.files['file']

    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(file.filename)[1]) as tmp:
            file.save(tmp)
            tmp_path = tmp.name

        extracted_text = extract_text(tmp_path)
        os.remove(tmp_path)

        if not extracted_text.strip():
            return jsonify({"error": "No readable text found."}), 400

        # Decide number of questions based on content length
        content_length = len(extracted_text)
        if content_length < 1000:
            num_questions = 10
        elif content_length < 3000:
            num_questions = 12
        elif content_length < 6000:
            num_questions = 15
        else:
            num_questions = 20

        # Limit context to avoid overloading the model
        combined_text = extracted_text[:43000]

        prompt = (
            "You are a professional academic quiz generator for college students. "
            "Based on the provided educational content, create a comprehensive quiz that covers all key concepts, facts, and definitions. "
            f"Generate {num_questions} questions (minimum 10, maximum 20, depending on content size). "
            "Use a mix of question types: multiple choice, true/false, and short answer. "
            "For multiple choice, provide 4 options (A–D) and indicate the correct answer letter. "
            "For true/false, provide the statement and the answer. "
            "For short answer, provide the question and the answer. "
            "Distribute the questions to cover as much of the material as possible, not just the beginning. "
            "Format the quiz as follows:\n\n"
            "Q: <question>\n"
            "Type: <Multiple Choice/True or False/Short Answer>\n"
            "A. <choice> (for MC only)\n"
            "B. <choice>\n"
            "C. <choice>\n"
            "D. <choice>\n"
            "Answer: <letter/True/False/short answer>\n"
            "----\n"
            "Content:\n"
            f"{combined_text}\n"
        )

        api_url = "https://openrouter.ai/api/v1/chat/completions"
        api_key = os.getenv("OPENROUTER_API_KEY")
        if not api_key:
            return jsonify({"error": "API key not set"}), 500

        headers = {
            "Authorization": f"Bearer {api_key}",
            "Content-Type": "application/json"
        }
        data = {
            "model": "openai/gpt-3.5-turbo",
            "messages": [{"role": "user", "content": prompt}],
            "temperature": 0.7
        }

        response = requests.post(api_url, headers=headers, json=data)
        result = response.json()
        if "choices" not in result or not result["choices"]:
            return jsonify({"error": result.get("error", "No 'choices' in API response"), "api_response": result}), 500

        content = result["choices"][0]["message"]["content"]

        # Extract questions using regex (handles all types)
        question_blocks = re.split(r'-{2,}', content)
        quiz = []
        for idx, block in enumerate(question_blocks, 1):
            q_match = re.search(r"Q:\s*(.*?)\nType:\s*(.*?)\n", block, re.DOTALL)
            if not q_match:
                continue
            question = q_match.group(1).strip()
            qtype = q_match.group(2).strip().lower()
            answer = ""
            choices = {}
            if "multiple" in qtype:
                a = re.search(r"A\.\s*(.*?)\n", block)
                b = re.search(r"B\.\s*(.*?)\n", block)
                c = re.search(r"C\.\s*(.*?)\n", block)
                d = re.search(r"D\.\s*(.*?)\n", block)
                answer_match = re.search(r"Answer:\s*([ABCD])", block)
                choices = {
                    "A": a.group(1).strip() if a else "",
                    "B": b.group(1).strip() if b else "",
                    "C": c.group(1).strip() if c else "",
                    "D": d.group(1).strip() if d else "",
                }
                answer = answer_match.group(1).strip() if answer_match else ""
            elif "true" in qtype:
                answer_match = re.search(r"Answer:\s*(True|False)", block, re.IGNORECASE)
                answer = answer_match.group(1).strip() if answer_match else ""
            elif "short" in qtype:
                answer_match = re.search(r"Answer:\s*(.*)", block)
                answer = answer_match.group(1).strip() if answer_match else ""
            quiz.append({
                "number": idx,
                "question": question,
                "type": qtype,
                "choices": choices,
                "answer": answer
            })

        # Only return questions that were parsed correctly
        quiz = [q for q in quiz if q["question"] and q["answer"]]

        return jsonify({"quiz": quiz})

    except Exception as e:
        return jsonify({"error": str(e)}), 500